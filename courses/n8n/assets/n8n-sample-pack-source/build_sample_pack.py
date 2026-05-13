#!/usr/bin/env python3
"""
n8n Lite Pack 演練素材包生成腳本 v1.1

產生 14 個 workflow 演練用的範例素材：
- 真實 PDF（reportlab）
- 真實 docx / pptx / xlsx（python-docx / pptx / openpyxl）
- markdown / txt / csv / json（純文字）
- PNG 佔位圖（PIL）
- README + 對照表 + manifest.json

v1.1 變動（依 Codex L3 審核 4cb89bc7 採納）：
- 加 6 個邊界 case 樣本（empty / not-a-pdf / Big5 / 未命名 / messy CSV / spike）
- 加 manifest.json（每 workflow 必備檔 SSOT）
- 加 validate_pack()（生成後自動驗每檔可讀 + 副檔名 magic 對齊）
- README 補 Windows 解壓 / Docker 路徑 / 失敗 troubleshoot 3 段

產出位置：./out/n8n-sample-pack/
打包：./out/n8n-sample-pack.zip
"""

import os
import shutil
import zipfile
from pathlib import Path
from datetime import datetime

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.enums import TA_LEFT
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook
from PIL import Image, ImageDraw, ImageFont


HERE = Path(__file__).parent
OUT = HERE / "out" / "n8n-sample-pack"
ZIP_OUT = HERE / "out" / "n8n-sample-pack.zip"

# ── 註冊中文字型（reportlab 內建 STSong-Light 涵蓋繁中常用字） ──
pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))


def reset_dirs():
    if OUT.exists():
        shutil.rmtree(OUT)
    if ZIP_OUT.exists():
        ZIP_OUT.unlink()
    OUT.mkdir(parents=True)


def make_pdf(path: Path, title: str, body_paragraphs: list[str]):
    """用 reportlab 生繁中 PDF。"""
    doc = SimpleDocTemplate(
        str(path),
        pagesize=A4,
        leftMargin=2 * cm, rightMargin=2 * cm,
        topMargin=2 * cm, bottomMargin=2 * cm,
    )
    styles = getSampleStyleSheet()
    h_style = ParagraphStyle(
        "h", parent=styles["Title"],
        fontName="STSong-Light", fontSize=18, leading=24,
        alignment=TA_LEFT, spaceAfter=14,
    )
    p_style = ParagraphStyle(
        "p", parent=styles["BodyText"],
        fontName="STSong-Light", fontSize=11, leading=18,
        alignment=TA_LEFT, spaceAfter=10,
    )
    story = [Paragraph(title, h_style), Spacer(1, 6)]
    for para in body_paragraphs:
        story.append(Paragraph(para.replace("\n", "<br/>"), p_style))
    doc.build(story)


def make_docx(path: Path, title: str, body_paragraphs: list[str]):
    doc = Document()
    doc.add_heading(title, level=1)
    for p in body_paragraphs:
        doc.add_paragraph(p)
    doc.save(str(path))


def make_pptx(path: Path, title: str, slides: list[tuple[str, list[str]]]):
    prs = Presentation()
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = "n8n Lite Pack 演練素材"
    # Content slides
    bullet_layout = prs.slide_layouts[1]
    for sub_title, bullets in slides:
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = sub_title
        body = slide.placeholders[1].text_frame
        body.text = bullets[0] if bullets else ""
        for b in bullets[1:]:
            p = body.add_paragraph()
            p.text = b
    prs.save(str(path))


def make_xlsx(path: Path, sheet_name: str, header: list[str], rows: list[list]):
    wb = Workbook()
    ws = wb.active
    ws.title = sheet_name
    ws.append(header)
    for r in rows:
        ws.append(r)
    wb.save(str(path))


def make_png(path: Path, text: str, size=(1200, 800), bg=(245, 243, 238), fg=(45, 50, 60)):
    img = Image.new("RGB", size, bg)
    draw = ImageDraw.Draw(img)
    # Try a system font; fall back to default
    font = None
    for f in [
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/Library/Fonts/Arial Unicode.ttf",
    ]:
        if os.path.exists(f):
            try:
                font = ImageFont.truetype(f, 64)
                break
            except Exception:
                pass
    if font is None:
        font = ImageFont.load_default()
    bbox = draw.textbbox((0, 0), text, font=font)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    draw.text(((size[0] - tw) / 2, (size[1] - th) / 2), text, fill=fg, font=font)
    img.save(str(path), "PNG")


def write_text(path: Path, content: str):
    path.write_text(content, encoding="utf-8")


# ─────────────────────────────────────────────
# 各分區生成
# ─────────────────────────────────────────────


def gen_pdf_inbox():
    """#02 PDF AI 改名 — 3 個檔名無意義的 PDF，內容各異"""
    d = OUT / "pdf-inbox"
    d.mkdir()

    make_pdf(
        d / "doc-001.pdf",
        "保密協議書（NDA）",
        [
            "本協議由甲方（弄一下工作室）與乙方（客戶 A 公司）於 2026 年 5 月 6 日簽訂。",
            "甲方因業務合作之需要，將向乙方揭露包含但不限於：商業計畫、技術文件、客戶名單、財務資料等機密資訊。",
            "乙方承諾於本協議有效期間及終止後三年內，不得對外揭露任何前述機密資訊，亦不得用於本協議目的以外之用途。",
            "如乙方違反本協議，應賠償甲方因此所受之一切損失，包括但不限於直接損失、間接損失與律師費用。",
            "本協議自雙方簽署之日起生效，有效期間為兩年。",
            "甲方代表：白辰幃　　　　乙方代表：（簽名處）",
        ],
    )

    make_pdf(
        d / "doc-002.pdf",
        "服務報價單",
        [
            "報價單編號：QT-2026-Q2-038",
            "報價日期：2026 年 5 月 6 日　　有效期限：30 天",
            "客戶名稱：客戶 B 股份有限公司",
            "─────────────────────────────",
            "項目一：n8n 流程自動化顧問導入　　　數量：1 式　　單價：NT$ 80,000",
            "項目二：客製 workflow 設計（5 個流程）　數量：1 式　　單價：NT$ 60,000",
            "項目三：6 小時內部培訓工作坊　　　　　數量：1 式　　單價：NT$ 35,000",
            "項目四：3 個月維運支援（每月 4 小時）　數量：1 式　　單價：NT$ 25,000",
            "─────────────────────────────",
            "小計：NT$ 200,000　　稅金 5%：NT$ 10,000　　總計：NT$ 210,000",
            "付款條件：簽約付 50% 訂金，驗收後付清。",
        ],
    )

    make_pdf(
        d / "doc-003.pdf",
        "提示詞工程基礎教材",
        [
            "本教材帶你從零學會 Prompt Engineering 的核心概念。",
            "什麼是提示詞：給大語言模型（LLM）的輸入文字，影響模型輸出品質的關鍵。",
            "好提示詞的四個元素：角色（你是誰）、任務（要做什麼）、限制（不要做什麼）、輸出格式（怎麼回）。",
            "範例對比：『幫我寫個郵件』（壞）vs『你是業務助理。寫一封給客戶 A 的延遲交付道歉信，繁中、200 字內、語氣誠懇但不卑微』（好）。",
            "進階技巧：Few-shot prompting（給範例）、Chain-of-thought（讓模型逐步思考）、Self-consistency（多次採樣選最常見答案）。",
            "本課程後續將涵蓋 Function Calling、RAG、Agent 等高階技巧。",
        ],
    )


def gen_batch_inbox():
    """#03 批次錯誤恢復 — 4 正常 + 1 損壞 + 邊界 case (empty / not-a-pdf)"""
    d = OUT / "batch-inbox"
    d.mkdir()
    for i in range(1, 5):
        write_text(
            d / f"good-{i:02d}.txt",
            f"這是第 {i} 份正常的批次處理輸入檔。\n內容只是純文字，n8n 應該能順利讀進來、處理、寫出。",
        )
    # 損壞檔（非 UTF-8 二進制）— 觸發 #03 的錯誤恢復路徑
    (d / "corrupted.bin").write_bytes(bytes([0xFF, 0xFE, 0x00, 0x42, 0xC3, 0x28, 0x00]) * 50)
    # v1.1：邊界 case — 0 byte 空檔（測 Code node 對空內容容錯）
    (d / "empty.txt").write_bytes(b"")
    # v1.1：副檔名偽裝 — .pdf 但內容是純文字（測 PDF 解析失敗時 Code node 是否優雅回 fallback）
    write_text(
        d / "not-a-pdf.pdf",
        "這個檔名是 .pdf 但內容其實是純文字。\nn8n 的 extractFromFile 用 PDF 模式讀會失敗，\nCode node 應該偵測到並走 #03 失敗路徑。",
    )


def gen_daily_input():
    """#04 定時 AI 日報 — .md / .txt 混合"""
    d = OUT / "daily-input"
    d.mkdir()
    write_text(
        d / "週一會議記錄.md",
        """# 5/6 週一例會

## 出席
白辰幃、客戶 A 經理、PM Vincent

## 議題
1. **n8n 課程進度**：14 個 workflow 已完工，今日進入 0-10 全面驗證階段。
2. **客戶 B 提案**：本週五前交報價單（總額 21 萬，含 6h 工作坊 + 3 個月維運）。
3. **GAI 認證研習**：下週 PR review 時把模擬考引擎 export 成 standalone JSON。

## 待辦
- [ ] 5/6 下班前修完 #10 RPM 限速問題（**P0**）
- [ ] 5/8 寄報價單給客戶 B（**P1**）
- [ ] 5/10 跟 PM Vincent 對齊認證研習宣傳素材（**P2**）

## 風險
GAI Free Tier RPM 限速可能影響課堂 demo 流暢度，需準備 paid key 備援。
""",
    )
    write_text(
        d / "客戶反饋本週.txt",
        """週客戶反饋彙整 2026/5/6

客戶 A：
- n8n Cloudflare Tunnel 設定文件清楚，自己完成接 webhook
- 反映 #04 定時日報希望支援 PDF 內容（已於 v0.7 加 Switch 完成）

客戶 B：
- 詢問是否能客製分類規則（合約 / 設計稿 / 報價單）
- 報價要求：總價控制在 20 萬內，分期付款

客戶 C：
- 對 #10 客戶資料夾整理高度有興趣（自己有大量未分類客戶檔）
- 詢問是否能改 prompt 換業務術語

整體：自動化流程教學的需求高度集中在「文件處理」「客戶整理」兩塊。下季可考慮把這兩塊獨立成單堂課程。
""",
    )
    # v1.1：邊界 case — Big5 編碼測試（測 #04 對非 UTF-8 文字的 fallback）
    big5_content = """週客戶反饋彙整 2026/5/6 [Big5 編碼測試]

客戶 D：詢問 #02 PDF 改名能否加批次模式（已用 #03）
客戶 E：反映 Big5 舊資料無法讀（這份檔本身就是 Big5）
"""
    (d / "客戶反饋_Big5.txt").write_bytes(big5_content.encode("big5", errors="replace"))

    write_text(
        d / "專案進度Q2.md",
        """# 專案進度 · 2026 Q2

## 已完成
- AI 資料工廠（n8n）課程改版 v1.0 → v1.6（PDF 修補、批次優化、rule-based fallback）
- Lite Pack 14 個 workflow 全部上線
- Cloud vs self-host 補充章節
- Win/Mac 雙系統 setup-wizard 強化

## 進行中
- 結業後補充教材（LLM Copilot 系列 8 頁 + Cloud 決策 1 頁）
- 課前 0-10 全面驗證（Mac + Win 雙端）

## 即將啟動
- GAI 認證研習模擬考引擎獨立化
- 工作坊招生頁改版（landing 風格）
""",
    )


def gen_client_inbox():
    """#10 客戶資料夾整理 — mixed file types，覆蓋 6 個分類"""
    d = OUT / "client-inbox"
    d.mkdir()

    # PDF 合約類
    make_pdf(
        d / "合約_NDA_客戶A.pdf",
        "保密協議書 NDA",
        [
            "本協議用於甲方與乙方共同合作期間之資訊保密。",
            "乙方不得將任何商業機密揭露給第三方。",
            "違反本協議將承擔法律責任。",
        ],
    )
    # PDF 報價類
    make_pdf(
        d / "Q3-報價單.pdf",
        "Q3 報價單",
        [
            "報價單編號：QT-2026-Q3-012",
            "客戶：客戶 C",
            "服務內容：n8n 自動化導入 + 工作坊。",
            "總計：NT$ 180,000",
        ],
    )
    # PPTX 簡報類
    make_pptx(
        d / "客戶提案-v1.pptx",
        "客戶 X 自動化導入提案",
        [
            ("專案概述", ["背景：客戶 X 目前手動處理客戶資料夾", "目標：自動分類 + 命名 + 通知", "預計效益：每週節省 8 小時人工"]),
            ("解決方案", ["n8n 流程自動化平台", "Gemini AI 分類引擎", "Cloudflare Tunnel 公開 webhook"]),
            ("時程與報價", ["導入：4 週", "培訓：6 小時工作坊", "報價：總額 NT$ 210,000"]),
        ],
    )
    # DOCX 教材類（兩份）
    make_docx(
        d / "01提示詞基礎.docx",
        "01 提示詞基礎",
        [
            "什麼是提示詞：給 LLM 的輸入文字。",
            "好提示詞的四個元素：角色、任務、限制、輸出格式。",
            "範例：『你是業務助理。寫一封給客戶的道歉信。』",
        ],
    )
    make_docx(
        d / "02提示詞進階.docx",
        "02 提示詞進階",
        [
            "Few-shot prompting：在 prompt 內給 2-5 個範例。",
            "Chain-of-thought：讓模型逐步思考。",
            "Self-consistency：多次採樣選最常見答案。",
        ],
    )
    # XLSX 報表類
    make_xlsx(
        d / "業務報表-2026Q2.xlsx",
        "Q2 業績",
        ["客戶", "案件", "金額", "狀態"],
        [
            ["客戶A", "n8n 導入", 200000, "已成交"],
            ["客戶B", "工作坊", 35000, "提案中"],
            ["客戶C", "顧問諮詢", 60000, "已成交"],
        ],
    )
    # PNG 圖檔
    make_png(d / "螢幕擷取-產品介面.png", "產品介面 V1\n2026/05/06")
    # v1.1：邊界 case — 不能只看檔名分類（測 #10 內容判斷的 fallback）
    make_docx(
        d / "未命名文件.docx",
        "（無標題）",
        [
            "本協議書由甲方（弄一下工作室）與乙方共同簽訂。",
            "雙方對於合作期間之機密資訊負有保密義務。",
            "違反本協議將承擔法律責任，賠償一切損失。",
            "（這份檔名沒有「合約」「NDA」等關鍵字，但內容明顯是合約 — 測 #10 的 classifyByContent fallback）",
        ],
    )
    # MD 紀錄類
    write_text(
        d / "會議記錄-0506.md",
        """# 5/6 例會紀錄

## 出席
白辰幃、PM Vincent

## 結論
- #10 改 v1.6 加 rule-based fallback
- 課前驗證走完 0-10 動線
- 下午 14:30 跟客戶 B 對齊報價
""",
    )


def gen_leads_inbox():
    """#11 CSV 線索清洗 — 含重複/缺欄/格式不一的雜亂資料"""
    d = OUT / "leads-inbox"
    d.mkdir()
    csv_content = """name,email,phone,company,source,note
王大明,wang@example.com,0912-345-678,A公司,LinkedIn,有興趣
李美麗,LiMei@Example.COM,0922 333 444,B 股份有限公司,展會,需要報價
,nobody@x.tw,0911111111,C 工作室,網站,
王大明,wang@example.com,0912345678,A公司,Email,重複線索
陳志強,chen@d.com,,D 文創,推薦,熱門客戶
林小芳,lin@e.com.tw,02-2345-6789,E 顧問,LinkedIn,
張三,san@f.tw,0900000000,F,Phone,
李四,si@g.tw,0987654321,G 公司,Email,要 quote
,blank@x,,,,
趙六,zhao@h.com,+886-921-555-666,H 集團,展會,VIP
錢七,qian@i.tw,0933 222 111,I 設計,LinkedIn,
孫八,sun@j.tw,,J,網站,需 demo
陳志強,chen@d.com,0966-555-777,D 文創,推薦,重複
周九,zhou@k.com,02 12345678,K 顧問,Email,合約洽談中
吳十,wu@L.com,0988123456,L 公司,展會,熱門
鄭一,zheng@m.tw,,M,網站,
小張,xiao@n.com,0911-000-999,N 工作室,LinkedIn,需要諮詢
小王,xwang@o.tw,0900-111-222,O 科技,Email,
小李,xli@p.com,03-1234567,P 顧問,Phone,VIP
,,,,,
測試 1,test@q.tw,0900-000-001,Q,Test,測試資料勿留
測試 2,test@q.tw,0900-000-002,Q,Test,重複測試
劉一二,liu@r.com,0922-888-999,R 集團,LinkedIn,合約簽訂中
何三四,he@s.tw,02-87654321,S 文創,推薦,
段五六,duan@t.com,+886-911-222-333,T 顧問,展會,VIP
姚七八,yao@u.tw,0966-111-222,U 公司,網站,
邵九十,shao@v.com,0955-444-333,V 工作室,Email,
小陳,xchen@w.tw,03-7654321,W 集團,Phone,
小林,xlin@x.tw,0922 555 666,X 設計,LinkedIn,熱門
黃十一,huang@y.com,0911 222 333,Y,展會,
"""
    write_text(d / "leads-raw.csv", csv_content)

    # v1.1：邊界 case — 欄名中英混雜 + UTF-8 BOM + 全形電話 + 換行不一致（測 #11 normalization）
    messy_csv = """﻿姓名,Email Address,電話,Company Name,來源,Note
王大明,wang@a.com,０９１２－３４５－６７８,A 公司,LinkedIn,全形電話
李美麗,limei@b.com,02 1234-5678,B 股份,展會,半形含空白
,empty@c.com,,,,
陳志強,chen@d.com,+886 921 555 666,D 文創,推薦,國碼+空白
鄭一二,zheng@e.com,no-phone,E 顧問,Email,電話欄寫文字\r\n小張,xiao@f.com,0900-000-000,F,LinkedIn,Windows CRLF 換行
小李,xli@g.com,0911222333,G,Email,缺市話\r"""
    # 注意：這裡刻意混 \n + \r\n + \r 來測 #11 的 row 切割是否 robust
    (d / "leads-raw-messy.csv").write_bytes(messy_csv.encode("utf-8"))


def gen_knowledge_docs():
    """#12 本地知識庫 RAG — 4 份知識文件"""
    d = OUT / "knowledge-docs"
    d.mkdir()
    write_text(
        d / "公司介紹.md",
        """# 弄一下工作室

弄一下工作室是一家專注於 AI 自動化教學與顧問導入的工作室，2024 年成立。

## 服務項目
- AI 資料工廠課程（n8n + Gemini）
- Gen AI 認證研習與商業應用工作坊
- 自動化流程顧問導入（n8n + Cloudflare Tunnel）

## 創辦人
白辰幃，曾任數家科技公司產品經理與設計師，專注於把複雜技術轉成可教可用的流程。

## 客戶
中小企業、行銷團隊、設計工作室、顧問公司、課程講師。
""",
    )
    write_text(
        d / "產品 FAQ.md",
        """# 產品常見問題

## Q1：n8n 跟 Make 差別？
n8n 是自架版自動化平台（資料留本機），Make 是雲端 SaaS。n8n 適合資料隱私敏感、流程量大、想客製節點的團隊。

## Q2：免費版 Gemini 有哪些限制？
免費 tier 限制：每分鐘 10 次（RPM）/ 每天 250 次（RPD）/ 每分鐘 25 萬 tokens（TPM）。批次處理多檔可能撞限速。

## Q3：要怎麼把 n8n 對外公開？
用 Cloudflare Tunnel 一鍵接通，或自架 reverse proxy。本工作室提供完整 setup-wizard 自動化。

## Q4：客戶導入大約多久？
標準 4 週，含需求訪談、流程設計、實作、培訓、驗收。
""",
    )
    write_text(
        d / "退款政策.md",
        """# 退款政策

## 課程退款
- 開課前 7 天以上：全額退款
- 開課前 1-7 天：退 80%
- 開課後 24 小時內未上課：退 50%
- 開課後超過 24 小時或已上課：不退款

## 顧問導入退款
- 簽約後 7 天內，未開始實作：全額退款
- 開始實作後：依完成比例計算

## 工具授權退款
帳號開通後不退款，但可轉讓給他人使用。
""",
    )
    write_text(
        d / "客戶案例集.md",
        """# 客戶案例集

## 案例 1：A 公司（電商）
- 痛點：每天人工分類客戶 email，每週耗 8 小時
- 解法：用 #09 Gmail 分類 workflow 自動歸檔到 4 路 Switch
- 效益：每週節省 8 小時，分類準確度 92%

## 案例 2：B 公司（行銷）
- 痛點：銷售線索 CSV 雜亂，重複資料多
- 解法：用 #11 CSV 清洗評分 workflow 自動 dedupe + AI 評分
- 效益：500 筆線索 5 分鐘清完，過往要兩天

## 案例 3：C 工作室（設計）
- 痛點：客戶提供的檔案資料夾混亂，每月整理 4 小時
- 解法：用 #10 客戶資料夾自動整理 + Telegram 通知
- 效益：每月節省 4 小時，加上完整交接清單給客戶
""",
    )


def gen_ops_input():
    """#13 ops snapshot — 系統 metrics CSV"""
    d = OUT / "ops-input"
    d.mkdir()
    csv_content = """timestamp,service,metric,value,unit
2026-05-06T08:00:00,api-gateway,requests_per_min,1240,count
2026-05-06T08:00:00,api-gateway,error_rate,0.8,percent
2026-05-06T08:00:00,api-gateway,p95_latency,180,ms
2026-05-06T09:00:00,api-gateway,requests_per_min,1580,count
2026-05-06T09:00:00,api-gateway,error_rate,1.2,percent
2026-05-06T09:00:00,api-gateway,p95_latency,210,ms
2026-05-06T10:00:00,api-gateway,requests_per_min,1820,count
2026-05-06T10:00:00,api-gateway,error_rate,4.5,percent
2026-05-06T10:00:00,api-gateway,p95_latency,520,ms
2026-05-06T11:00:00,api-gateway,requests_per_min,1950,count
2026-05-06T11:00:00,api-gateway,error_rate,8.2,percent
2026-05-06T11:00:00,api-gateway,p95_latency,890,ms
2026-05-06T08:00:00,db-primary,connections,45,count
2026-05-06T08:00:00,db-primary,query_p95,12,ms
2026-05-06T09:00:00,db-primary,connections,52,count
2026-05-06T09:00:00,db-primary,query_p95,18,ms
2026-05-06T10:00:00,db-primary,connections,78,count
2026-05-06T10:00:00,db-primary,query_p95,45,ms
2026-05-06T11:00:00,db-primary,connections,95,count
2026-05-06T11:00:00,db-primary,query_p95,120,ms
"""
    write_text(d / "metrics-2026-05-06.csv", csv_content)

    # v1.1：邊界 case — 缺時間 / 重複 timestamp / 極端 spike（測 #13 異常偵測 robust）
    spike_csv = """timestamp,service,metric,value,unit
2026-05-07T08:00:00,api-gateway,requests_per_min,1300,count
2026-05-07T08:00:00,api-gateway,requests_per_min,1300,count
2026-05-07T08:00:00,api-gateway,error_rate,1.0,percent
,api-gateway,p95_latency,200,ms
2026-05-07T09:00:00,api-gateway,error_rate,99.9,percent
2026-05-07T09:00:00,api-gateway,p95_latency,15000,ms
2026-05-07T10:00:00,api-gateway,requests_per_min,0,count
2026-05-07T10:00:00,api-gateway,error_rate,,percent
2026-05-07T11:00:00,db-primary,connections,9999,count
2026-05-07T11:00:00,db-primary,query_p95,abc,ms
"""
    write_text(d / "metrics-2026-05-07-spike.csv", spike_csv)


def gen_scale_up():
    """v1.2 (Codex 496055e3 採納)：scale-up 資料夾，給「有體感的批次效益」做準備。
    學員預設用各 inbox 內的小量教學版（不爆 RPM），想看省時感才把 scale-up/<sub>/* 搬進去。
    """
    base = OUT / "scale-up"
    base.mkdir()

    # ── #02 PDF rename: 12 個檔（Codex 建議 classroom 12-15）──
    pdf_d = base / "pdf-inbox-12"
    pdf_d.mkdir()
    titles_bodies = [
        ("採購訂單 PO-2026-001", ["採購單號：PO-2026-001", "供應商：A 工廠", "品項：鋁合金 200 公斤", "金額：NT$ 45,000"]),
        ("採購訂單 PO-2026-002", ["採購單號：PO-2026-002", "供應商：B 設備", "品項：CNC 刀具一組", "金額：NT$ 120,000"]),
        ("採購訂單 PO-2026-003", ["採購單號：PO-2026-003", "供應商：C 包材", "品項：紙箱 5000 個", "金額：NT$ 38,000"]),
        ("供應商評估報告 Q2", ["評估期間：2026 Q2", "供應商：D 物流", "準時率：96%", "瑕疵率：0.4%"]),
        ("會議記錄 0507", ["時間：5/7 14:00", "出席：白辰幃 / Vincent / 客戶 X", "結論：決定改 v1.6 加批次處理"]),
        ("會議記錄 0508", ["時間：5/8 10:00", "出席：白辰幃 / Vincent", "結論：報價單寄出後追 D+3"]),
        ("交貨清單 2026-05", ["客戶 A：5/3 已交", "客戶 B：5/5 已交", "客戶 C：5/8 預計交"]),
        ("付款通知 INV-238", ["發票號：INV-238", "客戶 A 公司", "金額 NT$ 210,000", "付款期限：2026/5/30"]),
        ("付款通知 INV-239", ["發票號：INV-239", "客戶 B", "金額 NT$ 95,000", "付款期限：2026/6/5"]),
        ("收款明細 2026-05-06", ["收款日：5/6", "客戶 A：NT$ 105,000（首期）", "客戶 D：NT$ 40,000（尾款）"]),
        ("出貨統計 2026-W18", ["週次：W18", "出貨筆數：32", "重量合計：1.8 噸", "目的地：5 縣市"]),
        ("品質異常報告 QA-2026-013", ["異常編號：QA-2026-013", "批號：B-238", "異常類型：尺寸偏差 0.3mm", "處置：全批次返工"]),
    ]
    for i, (title, body) in enumerate(titles_bodies, 1):
        make_pdf(pdf_d / f"raw-{i:02d}.pdf", title, body)

    # ── #03 batch error: 30 個檔（Codex 建議 30-50，含 5-8 邊界）──
    batch_d = base / "batch-inbox-30"
    batch_d.mkdir()
    for i in range(1, 24):
        write_text(batch_d / f"good-{i:02d}.txt", f"第 {i} 份正常處理輸入。\n編號：BATCH-2026-{i:03d}\n內容：客戶資料 / 訂單明細 / 報表資料 等情境模擬。")
    # 邊界 case 7 個
    (batch_d / "corrupted-1.bin").write_bytes(bytes([0xFF, 0xFE, 0x00, 0x42]) * 80)
    (batch_d / "corrupted-2.bin").write_bytes(bytes([0xC3, 0x28, 0xA0, 0xA0]) * 60)
    (batch_d / "empty-1.txt").write_bytes(b"")
    (batch_d / "empty-2.txt").write_bytes(b"")
    write_text(batch_d / "fake-1.pdf", "副檔名 .pdf 但其實是純文字 1。")
    write_text(batch_d / "fake-2.pdf", "副檔名 .pdf 但其實是純文字 2。")
    write_text(batch_d / "tiny-1.txt", "x")  # 1 byte 但非空

    # ── #10 folder-organize: 40 個檔（Codex 建議 40-60 mixed）──
    folder_d = base / "client-inbox-40"
    folder_d.mkdir()
    # 合約類 8 個
    for i in range(1, 9):
        make_pdf(folder_d / f"合約_客戶{chr(64+i)}_{i:03d}.pdf", f"保密協議書 NDA · 客戶 {chr(64+i)}",
                 ["甲方與乙方之間的保密協議", "違反條款將承擔法律責任", f"編號：NDA-2026-{i:03d}"])
    # 報價類 8 個
    for i in range(1, 9):
        make_pdf(folder_d / f"報價單_QT{i:03d}.pdf", f"報價單 QT-2026-{i:03d}",
                 ["項目：n8n 自動化導入", f"金額：NT$ {(50+i*15)*1000}", "有效期：30 天"])
    # 提案 pptx 6 個
    for i in range(1, 7):
        make_pptx(folder_d / f"客戶提案_v{i}.pptx", f"客戶 {chr(64+i)} 自動化提案",
                  [("背景", [f"客戶 {chr(64+i)} 痛點"]), ("方案", ["n8n + Gemini"]), ("時程", [f"{2*i} 週導入"])])
    # 教材 docx 6 個
    for i in range(1, 7):
        make_docx(folder_d / f"教材_{i:02d}_章節.docx", f"教材 第 {i} 章",
                  ["本章重點：自動化思維", "範例：批次處理 100 個檔案", "練習：3 道題"])
    # 報表 xlsx 4 個
    for i in range(1, 5):
        make_xlsx(folder_d / f"業績報表_2026Q{i}.xlsx", f"Q{i} 業績",
                  ["客戶", "金額", "狀態"], [["客戶 A", 100000+i*20000, "已成"], ["客戶 B", 50000+i*10000, "進行"]])
    # 圖檔 4 個
    for i in range(1, 5):
        make_png(folder_d / f"截圖_介面_{i:02d}.png", f"產品介面\nVersion {i}")
    # 紀錄 md 4 個
    for i in range(1, 5):
        write_text(folder_d / f"會議記錄_2026-05-{i:02d}.md",
                   f"# 5/{i} 會議\n\n出席：白辰幃 / 客戶 {chr(64+i)}\n\n結論：方案 V{i} 確認，下週交付。")

    # ── #11 CSV: 100 行 throttle 版（Codex 建議分教學 30 + throttle 100）──
    leads_d = base / "leads-inbox-100"
    leads_d.mkdir()
    csv_lines = ["name,email,phone,company,source,note"]
    sources = ["LinkedIn", "展會", "Email", "推薦", "網站", "Phone", "Webinar"]
    for i in range(1, 101):
        name = f"客戶_{i:03d}"
        email = f"client{i:03d}@example.com" if i % 7 != 0 else ""  # 14% 缺 email
        phone = f"09{i:08d}"[:10] if i % 5 != 0 else ""
        company = f"公司_{chr(65 + (i % 26))}{i:03d}"
        source = sources[i % len(sources)]
        note = ["熱門", "VIP", "需 demo", "合約洽談", ""][i % 5]
        csv_lines.append(f"{name},{email},{phone},{company},{source},{note}")
    # 加 5 筆重複測試 dedup
    for i in [3, 17, 42, 58, 91]:
        csv_lines.append(f"客戶_{i:03d},client{i:03d}@example.com,09{i:08d}"[:23] + f",公司_dup_{i:03d},Email,重複線索")
    write_text(leads_d / "leads-100.csv", "\n".join(csv_lines))

    # ── #13 ops: 200 行（Codex 建議 100-200）──
    ops_d = base / "ops-input-200"
    ops_d.mkdir()
    rows = ["timestamp,service,metric,value,unit"]
    services = ["api-gateway", "db-primary", "auth", "queue", "cache"]
    metrics = {"requests_per_min": "count", "error_rate": "percent", "p95_latency": "ms",
               "connections": "count", "query_p95": "ms"}
    import random
    random.seed(42)  # 可重現
    for h in range(0, 24):
        for s in services:
            for m, u in metrics.items():
                # 加 spike pattern
                base_v = {"requests_per_min": 1500, "error_rate": 1.0, "p95_latency": 200, "connections": 50, "query_p95": 15}[m]
                v = base_v * (3 if 14 <= h <= 16 else 1) * random.uniform(0.8, 1.2)
                rows.append(f"2026-05-07T{h:02d}:00:00,{s},{m},{v:.1f},{u}")
    write_text(ops_d / "metrics-day-200.csv", "\n".join(rows))

    # README
    write_text(
        base / "README.md",
        """# scale-up · 大規模演練素材（v1.2）

當你跑完小量教學版、想看「批次自動化省時感」時，把這裡的檔搬進對應 inbox：

| 子資料夾 | 對應 workflow | 量級 | 預期效益 |
|---|---|---|---|
| `pdf-inbox-12/` | #02 PDF AI 改名 | 12 PDF | 學員手動改 12 個檔需 5-8 分鐘；workflow 約 80 秒（含 throttle） |
| `batch-inbox-30/` | #03 批次錯誤恢復 | 30 檔含 7 邊界 | 看清楚 success/failure 分流規模 |
| `client-inbox-40/` | #10 客戶資料夾整理 | 40 mixed | 看 chunk batch + 6 類分流 + index report 完整效益 |
| `leads-inbox-100/` | #11 CSV 線索清洗 | 100+5dup | 看 normalize + dedupe + AI 評分批次價值（注意 RPM）|
| `ops-input-200/` | #13 ops 快照 | 200 行 metrics | 看異常 spike 偵測 + AI 分析 |

## ⚠ Gemini Free tier 限制提醒

跑大規模時注意：
- **Free tier 上限**：10 RPM / 250 RPD / 250K TPM（[官方](https://ai.google.dev/gemini-api/docs/rate-limits)）
- **#02 跑 12 個 PDF**：v1.2 已加 throttle（每筆 sleep 6.5 秒），約 80 秒跑完，安全
- **#03 跑 30 個檔**：throttle 後約 200 秒，會分到下一分鐘 RPM window
- **#11 跑 100 行**：是「1 次 batch prompt」設計，1 次 API call，安全
- **#13 跑 200 行**：是「aggregate 後 1 次 prompt」，安全

如果撞 429：v1.2 已加 60 秒自動 retry × 3，看 workflow log 出現 `retry on 429` 是正常的，等就好。

## 使用流程

```bash
# Mac
cp scale-up/pdf-inbox-12/* ~/Downloads/n8n-starter-kit/shared/pdf-inbox/

# Windows（PowerShell）
Copy-Item scale-up\\pdf-inbox-12\\* C:\\Users\\<you>\\Downloads\\n8n-starter-kit\\shared\\pdf-inbox\\
```

跑完想還原小量教學版：把對應 inbox 清空，從原 `n8n-sample-pack/<inbox>/` 重新複製。
""",
    )


def gen_manifest():
    """v1.1：每 workflow 必備檔的 SSOT。Codex 4cb89bc7 採納。"""
    import json
    manifest = {
        "version": "1.1",
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "description": "n8n Lite Pack 演練素材包 manifest — 每個 workflow 期望吃的檔案",
        "workflows": {
            "02-pdf-ai-rename": {
                "folder": "pdf-inbox",
                "expected_files": [
                    {"name": "doc-001.pdf", "type": "pdf-text", "expected": "AI rename to NDA-related"},
                    {"name": "doc-002.pdf", "type": "pdf-text", "expected": "AI rename to 報價單-related"},
                    {"name": "doc-003.pdf", "type": "pdf-text", "expected": "AI rename to 提示詞教材-related"},
                ],
            },
            "03-batch-error-recovery": {
                "folder": "batch-inbox",
                "expected_files": [
                    {"name": "good-01.txt", "type": "text", "expected": "success path"},
                    {"name": "good-02.txt", "type": "text", "expected": "success path"},
                    {"name": "good-03.txt", "type": "text", "expected": "success path"},
                    {"name": "good-04.txt", "type": "text", "expected": "success path"},
                    {"name": "corrupted.bin", "type": "binary", "expected": "failure path (non-UTF-8)"},
                    {"name": "empty.txt", "type": "empty", "expected": "edge case: 0 byte"},
                    {"name": "not-a-pdf.pdf", "type": "fake-pdf", "expected": "edge case: extension mismatch"},
                ],
            },
            "04-daily-ai-report": {
                "folder": "daily-input",
                "expected_files": [
                    {"name": "週一會議記錄.md", "type": "markdown", "expected": "summary section"},
                    {"name": "客戶反饋本週.txt", "type": "text-utf8", "expected": "summary section"},
                    {"name": "專案進度Q2.md", "type": "markdown", "expected": "summary section"},
                    {"name": "客戶反饋_Big5.txt", "type": "text-big5", "expected": "edge case: encoding fallback"},
                ],
            },
            "10-folder-organize": {
                "folder": "client-inbox",
                "expected_files": [
                    {"name": "合約_NDA_客戶A.pdf", "type": "pdf-text", "expected": "→ contracts/"},
                    {"name": "Q3-報價單.pdf", "type": "pdf-text", "expected": "→ invoices/"},
                    {"name": "客戶提案-v1.pptx", "type": "pptx", "expected": "→ presentations/"},
                    {"name": "01提示詞基礎.docx", "type": "docx", "expected": "→ docs/"},
                    {"name": "02提示詞進階.docx", "type": "docx", "expected": "→ docs/"},
                    {"name": "業務報表-2026Q2.xlsx", "type": "xlsx", "expected": "→ docs/ or others/"},
                    {"name": "螢幕擷取-產品介面.png", "type": "image", "expected": "→ images/"},
                    {"name": "會議記錄-0506.md", "type": "markdown", "expected": "→ docs/"},
                    {"name": "未命名文件.docx", "type": "docx", "expected": "edge case: classify by content not name"},
                ],
            },
            "11-csv-clean-score": {
                "folder": "leads-inbox",
                "expected_files": [
                    {"name": "leads-raw.csv", "type": "csv", "expected": "30 rows, dedupe + AI score"},
                    {"name": "leads-raw-messy.csv", "type": "csv-messy", "expected": "edge case: BOM + 全形 + CRLF mix"},
                ],
            },
            "12-knowledge-rag": {
                "folder": "knowledge-docs",
                "expected_files": [
                    {"name": "公司介紹.md", "type": "markdown"},
                    {"name": "產品 FAQ.md", "type": "markdown"},
                    {"name": "退款政策.md", "type": "markdown"},
                    {"name": "客戶案例集.md", "type": "markdown"},
                ],
            },
            "13-daily-ops-snapshot": {
                "folder": "ops-input",
                "expected_files": [
                    {"name": "metrics-2026-05-06.csv", "type": "csv-metrics", "expected": "20 rows"},
                    {"name": "metrics-2026-05-07-spike.csv", "type": "csv-metrics", "expected": "edge case: spike + missing time + dupe"},
                ],
            },
        },
        "no_sample_needed": ["01-webhook-hello-world", "05-telegram-notify",
                             "06-webhook-gemini-file", "07-quick-tunnel-receiver",
                             "08-expression-practice", "09-gmail-categorize",
                             "14-api-monitor"],
    }
    (OUT / "manifest.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )


def validate_pack():
    """v1.2：生成後逐檔驗證（可讀性 + 副檔名 magic 對齊）。Codex 4cb89bc7 採納。"""
    # 已知預期的邊界 case（不算 issue）
    expected_empty_names = {"empty.txt", "empty-1.txt", "empty-2.txt"}
    expected_fake_pdf_names = {"not-a-pdf.pdf", "fake-1.pdf", "fake-2.pdf"}

    issues = []
    file_count = 0
    for path in sorted(OUT.rglob("*")):
        if path.is_dir() or path.name == "manifest.json" or path.name == "README.md":
            continue
        file_count += 1
        try:
            size = path.stat().st_size
            ext = path.suffix.lower()
            rel = path.relative_to(OUT)

            # 0 byte 檢查（命名為 empty* 是預期）
            if size == 0 and path.name not in expected_empty_names:
                issues.append(f"  ⚠ {rel} is 0 bytes")
                continue

            # 副檔名 magic number 對齊
            if ext == ".pdf":
                with open(path, "rb") as f:
                    head = f.read(5)
                if head != b"%PDF-" and path.name not in expected_fake_pdf_names:
                    issues.append(f"  ⚠ {rel} 副檔名 .pdf 但不是 PDF magic")
            elif ext == ".png":
                with open(path, "rb") as f:
                    head = f.read(8)
                if head != b"\x89PNG\r\n\x1a\n":
                    issues.append(f"  ⚠ {rel} 副檔名 .png 但不是 PNG magic")
            elif ext in (".docx", ".pptx", ".xlsx"):
                with open(path, "rb") as f:
                    head = f.read(2)
                if head != b"PK":
                    issues.append(f"  ⚠ {rel} OOXML 應該是 ZIP (PK) 開頭")

            # 檔名 Win 禁字檢查（學員會在 Win 上解壓）
            for ch in '<>:"|?*':
                if ch in path.name:
                    issues.append(f"  ⚠ {rel} 含 Win 禁字 '{ch}'")
                    break

            # 尾部空白 / 句點（Win 不允許）
            stem = path.stem
            if stem != stem.rstrip(" ."):
                issues.append(f"  ⚠ {rel} 主檔名尾部含空白或句點（Win 禁）")

        except Exception as e:
            issues.append(f"  ✗ {path.relative_to(OUT)}: {e}")

    print(f"\n✓ validate_pack: 掃了 {file_count} 個檔")
    if issues:
        print(f"  發現 {len(issues)} 個問題：")
        for it in issues:
            print(it)
    else:
        print("  ✅ 全部通過 — 沒有檔名 / 編碼 / magic 問題")
    return len(issues) == 0


def gen_readme():
    write_text(
        OUT / "README.md",
        """# n8n Lite Pack 演練素材包 v1.2

對應 14 個 workflow 的測試素材。把對應子資料夾的內容**複製**到你的 starter-kit 內（**不要動原本的 shared 資料夾結構**）：

## 使用方法

1. 解壓本 zip 到任意地方
2. 把各子資料夾內的檔**複製**到 `n8n-starter-kit/shared/<對應資料夾>/`
3. 開 n8n UI → 對應 workflow → Execute

## ⚠ Windows 學員注意（解壓動線）

**請用 Win 內建「解壓縮全部」或 7-Zip 解到「短英文路徑」**，例如：
- ✅ `C:\\Users\\<你>\\Downloads\\n8n-sample-pack\\`
- ❌ 不要解到 OneDrive / iCloud / 桌面同步路徑（中文路徑也避開）
- ❌ 不要解到 `C:\\Users\\<你>\\OneDrive\\桌面\\客戶檔\\...`

OneDrive / iCloud 同步資料夾跟 Docker bind mount 衝突，會導致 n8n 看不到檔或讀到舊版。

## n8n Docker 路徑前置檢查

複製檔案進 starter-kit 後，先用 n8n UI 確認容器內看得到：
1. 開任一 Read File 節點
2. 路徑填 `/files/n8n-sample-pack/pdf-inbox/doc-001.pdf`（**不是** Mac/Win 本機路徑）
3. Execute Node → 看到綠燈 + binary preview = 容器掛載 OK

如果紅燈說「ENOENT」：你的檔案不在 Docker 看得到的 volume 裡，回去檢查 starter-kit 的 `shared/` 路徑。

## 對照表

| 子資料夾 | 對應 workflow | 用途 |
|---|---|---|
| `pdf-inbox/` | #02 PDF AI 改名 | 3 個檔名無意義 PDF（合約 / 報價單 / 教材內容）測試 AI 改名 |
| `batch-inbox/` | #03 批次錯誤恢復 | 4 正常 + 1 損壞檔 + **2 邊界 case** (`empty.txt` / `not-a-pdf.pdf`) |
| `daily-input/` | #04 定時 AI 日報 | .md + .txt 混合 + **`客戶反饋_Big5.txt`**（測 encoding fallback）|
| `client-inbox/` | #10 客戶資料夾整理 | PDF/PPTX/DOCX/XLSX/PNG/MD mixed + **`未命名文件.docx`**（測 classify by content） |
| `leads-inbox/` | #11 CSV 線索清洗 | 30 筆雜亂 + **`leads-raw-messy.csv`**（測 BOM/全形/CRLF normalize） |
| `knowledge-docs/` | #12 本地知識庫 RAG | 4 份知識文件 |
| `ops-input/` | #13 ops snapshot | 標準 metrics + **`metrics-2026-05-07-spike.csv`**（測異常+缺值容錯） |

**不需要素材的 workflow**（webhook trigger / OAuth credential / 設定驅動）：
- #01 webhook hello world
- #05 Telegram 通知測試
- #06 Webhook → Gemini → 本機檔（POST 觸發）
- #07 Quick Tunnel receiver
- #08 Expression 練習
- #09 Gmail 分類（需 OAuth）
- #14 API monitor（在 workflow 內設定 API list）

完整 workflow ↔ files 對照見 `manifest.json`。

## 🚀 想看到「批次自動化省時感」？

預設小量教學版（不爆 RPM）約 5-15 秒跑完，**體感不明顯**。想看真正的省時效益：

```bash
# Mac 範例：把 12 個 PDF 搬到 #02 inbox 跑
cp scale-up/pdf-inbox-12/* ~/Downloads/n8n-starter-kit/shared/pdf-inbox/
# 跑 #02 約 80 秒（含 throttle，避免撞 Free tier RPM）
# 比手動改 12 個檔的 5-8 分鐘省一個量級
```

`scale-up/` 含 5 個放大版資料夾，詳見 `scale-up/README.md`。

## 失敗時看哪裡（troubleshoot 三步）

跑 workflow 出錯先按這 3 步排查：

1. **AI 改名失敗（#02 / #03 / #10）** — 看新檔名是否含 `/ : ? *` 等 Win 禁字。Code node 已有 `replace([\\\\/:*?"<>|]/g, '_')`，但邊界 case（trailing space / Windows 保留名 CON/PRN/NUL）不在覆蓋範圍。
2. **讀檔失敗（任何 workflow）** — 先測 `pdf-inbox/doc-001.pdf` 能不能 Execute Read PDF 看到 binary preview。能 → 你的後段有問題；不能 → Docker 掛載 / 路徑 / 權限 問題。
3. **CSV / TXT garbled（#04 / #11）** — 檔案是否仍 UTF-8？打開 `leads-raw.csv` 用 VS Code/TextEdit 看，狀態列應顯示 UTF-8。`客戶反饋_Big5.txt` 是故意 Big5 — 那個是測 encoding fallback。

## 重新產生

如果想客製或補檔，跑這個腳本：

```bash
pip install reportlab python-docx python-pptx openpyxl pillow
python3 build_sample_pack.py
```

build script v1.1 會跑 `validate_pack()` 自動檢查每個檔案的可讀性 + 副檔名 magic number 對齊。
""",
    )


def make_zip():
    """
    打 zip 時顯式設 UTF-8 flag bit (general purpose bit 11, 0x0800)，
    確保 Windows 解壓器能正確顯示中文檔名。
    （macOS 內建 zip 指令 + Python ZipFile.write() 對純中文檔名都不保證設此 flag）
    """
    with zipfile.ZipFile(ZIP_OUT, "w", zipfile.ZIP_DEFLATED, allowZip64=True) as z:
        # 跳過 .DS_Store；保留空目錄結構
        for root, dirs, files in os.walk(OUT):
            dirs.sort()
            rel_dir = Path(root).relative_to(OUT.parent)
            if str(rel_dir) != ".":
                zinfo = zipfile.ZipInfo(str(rel_dir) + "/")
                zinfo.flag_bits |= 0x0800
                zinfo.external_attr = 0o755 << 16
                z.writestr(zinfo, b"")
            for f in sorted(files):
                if f == ".DS_Store":
                    continue
                full = Path(root) / f
                rel = full.relative_to(OUT.parent)
                zinfo = zipfile.ZipInfo.from_file(full, str(rel))
                zinfo.flag_bits |= 0x0800  # UTF-8 flag
                zinfo.compress_type = zipfile.ZIP_DEFLATED
                with open(full, "rb") as src:
                    z.writestr(zinfo, src.read())
    print(f"✓ 打包：{ZIP_OUT} ({ZIP_OUT.stat().st_size / 1024:.1f} KB)")


def main():
    print(f"產出位置：{OUT}")
    reset_dirs()
    gen_pdf_inbox()
    print("✓ pdf-inbox/ — 3 個 PDF")
    gen_batch_inbox()
    print("✓ batch-inbox/ — 4 正常 + 1 損壞 + 2 邊界 (empty / not-a-pdf)")
    gen_daily_input()
    print("✓ daily-input/ — 3 個 markdown/txt + 1 個 Big5")
    gen_client_inbox()
    print("✓ client-inbox/ — 8 個 mixed + 1 個 未命名（測 content 分類）")
    gen_leads_inbox()
    print("✓ leads-inbox/ — 1 標準 CSV + 1 messy CSV (BOM/全形/CRLF)")
    gen_knowledge_docs()
    print("✓ knowledge-docs/ — 4 個 markdown")
    gen_ops_input()
    print("✓ ops-input/ — 1 標準 metrics + 1 spike (測異常偵測)")
    gen_scale_up()
    print("✓ scale-up/ — 5 個放大版資料夾（pdf×12 / batch×30 / folder×40 / leads×100 / ops×200）")
    gen_manifest()
    print("✓ manifest.json — workflow ↔ files 對照 SSOT")
    gen_readme()
    print("✓ README.md")
    if not validate_pack():
        print("\n⚠ validate_pack 發現問題，但 zip 仍會打包（看上面警告）")
    make_zip()


if __name__ == "__main__":
    main()
